#!/usr/bin/env python3
"""NES Presentation — FINAL Quality Review.
Quality over Quantity. One engineering question per slide. One strong visual. One key takeaway."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BG_DARK   = RGBColor(0x0B, 0x1D, 0x3A)
BG_MID    = RGBColor(0x12, 0x2B, 0x4F)
NAVY      = RGBColor(0x1B, 0x26, 0x3B)
BLUE      = RGBColor(0x00, 0x96, 0xD6)
LIGHT_BLUE= RGBColor(0x1A, 0x3A, 0x5C)
GREEN     = RGBColor(0x27, 0xAE, 0x60)
RED       = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE    = RGBColor(0xF3, 0x9C, 0x12)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x8E, 0x9E, 0xAB)

LEFT_M = Inches(0.8)
CHART  = "/Users/yashitarora/Desktop/NES/Nes1/charts"

def bg(s):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG_DARK

def add_shape(s, st, x, y, w, h, fill=None, line_color=None, line_width=None):
    sh = s.shapes.add_shape(st, x, y, w, h)
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else: sh.fill.background()
    if line_color: sh.line.color.rgb = line_color; sh.line.width = line_width or Pt(1)
    else: sh.line.fill.background()
    return sh

def txt(s, x, y, w, h, t, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    tx = s.shapes.add_textbox(x, y, w, h); tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = a
    r = p.add_run(); r.text = t; r.font.size = Pt(sz)
    r.font.color.rgb = c; r.font.bold = b; r.font.name = "Calibri"

def ml(s, x, y, w, h, items, sz=13):
    tx = s.shapes.add_textbox(x, y, w, h); tf = tx.text_frame; tf.word_wrap = True
    for i, (t, c, b) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = t; r.font.size = Pt(sz)
        r.font.color.rgb = c; r.font.bold = b; r.font.name = "Calibri"
        p.space_after = Pt(2)

def header_bar(s, title):
    add_shape(s, MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9), fill=NAVY)
    txt(s, LEFT_M, Inches(0.2), Inches(10), Inches(0.5), title, sz=26, c=WHITE, b=True)
    add_shape(s, MSO_SHAPE.RECTANGLE, LEFT_M, Inches(0.75), Inches(2), Pt(3), fill=BLUE)

def footer(s, n):
    add_shape(s, MSO_SHAPE.RECTANGLE, LEFT_M, Inches(7.0), Inches(11.7), Pt(1), fill=BLUE)
    txt(s, LEFT_M, Inches(7.1), Inches(5), Inches(0.3), "NES Technologies — DLMS Meter Data Analytics", sz=10, c=GRAY)
    txt(s, Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.3), str(n), sz=10, c=GRAY, a=PP_ALIGN.RIGHT)

def sh(s, y, text, x=None):
    x = x or LEFT_M
    add_shape(s, MSO_SHAPE.RECTANGLE, x, y, Inches(0.5), Pt(3), fill=BLUE)
    txt(s, x + Inches(0.7), y - Inches(0.05), Inches(2.0), Inches(0.35), text, sz=13, c=WHITE, b=True)

def tbl(s, rows, x, y, w, cw, rh=Inches(0.28)):
    nr, nc = len(rows), len(cw)
    ts = s.shapes.add_table(nr, nc, x, y, w, rh * nr); t = ts.table
    for ci, c in enumerate(cw): t.columns[ci].width = c
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci); cell.text = str(val); cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                if ri == 0: cell.fill.solid(); cell.fill.fore_color.rgb = NAVY; fc, fb = WHITE, True
                else: cell.fill.solid(); cell.fill.fore_color.rgb = BG_MID if ri % 2 == 1 else NAVY; fc, fb = WHITE, False
                for run in p.runs: run.font.size = Pt(12); run.font.name = "Calibri"; run.font.color.rgb = fc; run.font.bold = fb
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)

def obs(s, x, y, w, h, title, text, ac=BLUE):
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=LIGHT_BLUE)
    add_shape(s, MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), h, fill=ac)
    txt(s, x + Inches(0.15), y + Inches(0.03), w - Inches(0.3), Inches(0.2), title, sz=10, c=ac, b=True)
    txt(s, x + Inches(0.15), y + Inches(0.23), w - Inches(0.3), h - Inches(0.26), text, sz=10, c=WHITE)

def chrt(s, fn, x, y, w, h):
    p = os.path.join(CHART, fn)
    if os.path.exists(p): s.shapes.add_picture(p, x, y, w, h)

def kpi(s, x, y, w, h, title, value, desc, ac=BLUE):
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=BG_MID)
    add_shape(s, MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), h, fill=ac)
    txt(s, x + Inches(0.12), y + Inches(0.05), w - Inches(0.15), Inches(0.2), title, sz=9, c=GRAY)
    txt(s, x + Inches(0.12), y + Inches(0.25), w - Inches(0.15), Inches(0.35), value, sz=22, c=ac, b=True)
    txt(s, x + Inches(0.12), y + Inches(0.6), w - Inches(0.15), Inches(0.2), desc, sz=9, c=WHITE)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE (No KPIs — problem and approach only)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
add_shape(s, MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08), fill=BLUE)
logo = "/Users/yashitarora/Desktop/NES/Nes1/nes_logo_transparent.png"
if os.path.exists(logo): s.shapes.add_picture(logo, Inches(10.8), Inches(0.3), Inches(2.2), Inches(0.6))

txt(s, LEFT_M, Inches(1.0), Inches(10), Inches(0.6), "DLMS Smart Meter Data Analysis", sz=34, c=WHITE, b=True)
txt(s, LEFT_M, Inches(1.6), Inches(10), Inches(0.5), "Forecasting & Anomaly Detection for Power Distribution", sz=28, c=BLUE, b=True)
txt(s, LEFT_M, Inches(2.3), Inches(10), Inches(0.4), "UC2: Behaviour Analysis → Forecasting → Anomaly Detection  |  UC1: Investigation → Root Cause", sz=14, c=WHITE)
add_shape(s, MSO_SHAPE.RECTANGLE, LEFT_M, Inches(2.9), Inches(4), Pt(2), fill=BLUE)

txt(s, LEFT_M, Inches(3.2), Inches(10), Inches(0.35), "Business Problem", sz=14, c=WHITE, b=True)
txt(s, LEFT_M, Inches(3.6), Inches(10), Inches(0.5), "NES needs systematic detection of abnormal consumption patterns\nto reduce distribution losses and identify metering issues.", sz=14, c=WHITE)

txt(s, LEFT_M, Inches(4.5), Inches(10), Inches(0.35), "Scope", sz=14, c=WHITE, b=True)
txt(s, LEFT_M, Inches(4.9), Inches(10), Inches(0.5), "5 meters × 12 months × 175,200 readings  |  LightGBM + 6-algorithm ensemble", sz=14, c=WHITE)

txt(s, LEFT_M, Inches(5.8), Inches(11), Inches(0.25), "ENGINEERING PIPELINE", sz=12, c=WHITE, b=True)
steps = [("Data\nQuality", NAVY), ("Per-Meter\nBehaviour", GREEN), ("Monthly\nAnalysis", NAVY),
         ("Feature\nEng.", NAVY), ("Model\nSelection", GREEN), ("Forecast\nValidation", GREEN),
         ("UC2→UC1\nHandoff", ORANGE), ("UC1\nInvestigation", RED), ("Root\nCause", RED), ("Business\nRecommendation", NAVY)]
for i, (step, bc) in enumerate(steps):
    x = LEFT_M + Inches(i * 1.22)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(6.1), Inches(1.1), Inches(0.45), fill=bc)
    txt(s, x, Inches(6.12), Inches(1.1), Inches(0.42), step, sz=9, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        txt(s, x + Inches(1.12), Inches(6.15), Inches(0.1), Inches(0.35), ">", sz=10, c=GRAY, a=PP_ALIGN.CENTER)

footer(s, 1)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — DATA QUALITY (1 chart + 1 table)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header_bar(s, "Data Quality Assessment — DLMS Meter Artifacts")

sh(s, Inches(1.5), "REVERSE ENERGY — DLMS Firmware Artifact")
chrt(s, "reverse_energy_trend.png", LEFT_M, Inches(1.8), Inches(9.0), Inches(3.5))
obs(s, LEFT_M, Inches(5.4), Inches(9.0), Inches(0.4),
    "Finding", "~50% negative energy diffs across all meters = DLMS firmware artifact. No zero-voltage bypass detected.", GREEN)

sh(s, Inches(1.5), "QUALITY ISSUES PER METER", Inches(10.2))
tbl(s, [["Meter", "Issues", "Type"],
    ["MTR001", "17,590", "Negative currents"],
    ["MTR002", "17,476", "Reverse energy"],
    ["MTR003", "17,380", "Reverse energy"],
    ["MTR004", "17,498", "Reverse energy"],
    ["MTR005", "17,443", "Reverse energy"]],
    Inches(10.2), Inches(1.8), Inches(2.3), [Inches(0.5), Inches(0.7), Inches(1.1)], rh=Inches(0.24))

obs(s, Inches(10.2), Inches(4.0), Inches(2.3), Inches(0.4),
    "Total", "87,387 issues. All DLMS artifacts.", GRAY)

footer(s, 2)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — PER-METER BEHAVIOUR (1 chart + 1 table)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header_bar(s, "Per-Meter Behaviour — Each Meter is Different")

sh(s, Inches(1.5), "DAILY LOAD PROFILE — Per Meter (Full Width)")
chrt(s, "per_meter_daily_profile.png", LEFT_M, Inches(1.8), Inches(11.7), Inches(4.0))
obs(s, LEFT_M, Inches(5.9), Inches(11.7), Inches(0.4),
    "Key Insight", "MTR001 peak=3.3x trough (residential). MTR005 peak=1.8x (industrial 24h). Unique daily shapes justify per-meter modelling.", BLUE)

sh(s, Inches(6.4), "CONSUMER CLASSIFICATION — 4 Evidence Dimensions")
tbl(s, [["Meter", "Class", "Base kW", "Night/Day", "Load Factor"],
    ["MTR001", "Residential", "0.993", "0.40", "0.472"],
    ["MTR002", "Small Commercial", "1.486", "0.49", "0.498"],
    ["MTR003", "Workshop", "2.000", "0.55", "0.543"],
    ["MTR004", "Office", "2.507", "0.60", "0.581"],
    ["MTR005", "Industrial", "3.021", "0.64", "0.609"]],
    LEFT_M, Inches(6.7), Inches(11.7), [Inches(0.7), Inches(1.3), Inches(0.7), Inches(0.7), Inches(0.7)], rh=Inches(0.14))

footer(s, 3)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — MONTHLY ANALYSIS (1 chart + 1 table)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header_bar(s, "Monthly Analysis — Unique Per-Meter Fingerprint")

sh(s, Inches(1.5), "MONTHLY PEAK TIMELINE — Each Meter Has Distinct Peak Month")
chrt(s, "monthly_peak_timeline.png", LEFT_M, Inches(1.8), Inches(11.7), Inches(4.0))
obs(s, LEFT_M, Inches(5.9), Inches(11.7), Inches(0.4),
    "Key Insight", "MTR005 Nov=6.995kW (global peak). MTR004 Apr=6.445kW. Per-meter forecasting required — do not combine.", BLUE)

sh(s, Inches(6.4), "PEAK LOAD COMPARISON")
tbl(s, [["Meter", "Min Peak", "Max Peak", "Peak Month", "Avg kW"],
    ["MTR001", "4.380", "4.689", "Sep", "2.215"],
    ["MTR002", "4.963", "5.468", "Jul", "2.720"],
    ["MTR003", "5.371", "5.948", "Jul", "3.227"],
    ["MTR004", "5.935", "6.445", "Apr", "3.741"],
    ["MTR005", "6.548", "6.995", "Nov", "4.252"]],
    LEFT_M, Inches(6.7), Inches(11.7), [Inches(0.7), Inches(0.8), Inches(0.8), Inches(0.8), Inches(0.7)], rh=Inches(0.14))

footer(s, 4)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — FEATURE ENGINEERING (1 chart + 1 table)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header_bar(s, "Feature Engineering — What Drives the Model")

sh(s, Inches(1.5), "FEATURE IMPORTANCE — LightGBM Internal")
chrt(s, "feature_importance_lgbm.png", LEFT_M, Inches(1.8), Inches(10.0), Inches(4.0))
obs(s, LEFT_M, Inches(5.9), Inches(10.0), Inches(0.4),
    "Key Finding", "Lag features dominate. Interaction features (+29% impact when removed) are most critical.", GREEN)

sh(s, Inches(1.5), "FEATURE GROUPS — Impact When Removed", Inches(10.8))
tbl(s, [["Group", "Impact"],
    ["Interaction", "+29%"],
    ["Lag", "+4%"],
    ["Weather", "+1%"],
    ["Calendar", "+0%"]],
    Inches(10.8), Inches(1.8), Inches(1.7), [Inches(0.9), Inches(0.8)], rh=Inches(0.22))

footer(s, 5)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — MODEL SELECTION (NEW: Why LightGBM?)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header_bar(s, "Model Selection — Why LightGBM?")

sh(s, Inches(1.5), "MODEL COMPARISON — 3 Algorithms Evaluated (Corrected)")
tbl(s, [["Model", "MAE", "RMSE", "MAPE", "R²", "Verdict"],
    ["LightGBM", "0.0171", "0.0221", "0.72%", "0.9997", "SELECTED"],
    ["Linear Regression", "0.0252", "0.0380", "1.18%", "0.9990", "Good baseline"],
    ["Random Forest", "0.0495", "0.0641", "3.34%", "0.9971", "Slow, higher error"]],
    LEFT_M, Inches(1.8), Inches(11.7), [Inches(1.3), Inches(0.8), Inches(0.8), Inches(0.7), Inches(0.7), Inches(1.3)], rh=Inches(0.22))
obs(s, LEFT_M, Inches(3.2), Inches(11.7), Inches(0.4),
    "Decision", "LightGBM selected: MAPE=0.72%, R²=0.9997. 61% better than Linear Regression, 79% better than Random Forest.", GREEN)

sh(s, Inches(3.8), "LIGHTGBM HYPERPARAMETERS — As Configured in Notebook")
tbl(s, [["Parameter", "Value", "Reason"],
    ["n_estimators", "300", "Early stopping at 50 rounds"],
    ["learning_rate", "0.1", "Default, works well"],
    ["num_leaves", "31", "Default, controls complexity"],
    ["min_child_samples", "20", "Smooths predictions"],
    ["feature_fraction", "0.8", "Reduces overfitting"]],
    LEFT_M, Inches(4.1), Inches(11.7), [Inches(1.5), Inches(1.0), Inches(3.5)], rh=Inches(0.20))

footer(s, 6)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — FORECAST VALIDATION (1 chart + 1 table)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header_bar(s, "Forecast Validation — MAPE = 0.72%, R² = 0.9997")

sh(s, Inches(1.5), "ACTUAL vs PREDICTED — 14,400 Test Readings (Full Width)")
chrt(s, "forecast_validation.png", LEFT_M, Inches(1.8), Inches(11.7), Inches(4.0))
obs(s, LEFT_M, Inches(5.9), Inches(11.7), Inches(0.4),
    "Proof", "R²=0.9997. Residuals centered at 0. No systematic bias. Model is unbiased and accurate.", GREEN)

sh(s, Inches(6.4), "PER-METER ACCURACY")
tbl(s, [["Meter", "MAE", "MAPE%", "R²", "Status"],
    ["MTR001", "0.017", "1.42%", "0.9994", "Highest error (residential)"],
    ["MTR002", "0.017", "0.72%", "0.9995", "Good"],
    ["MTR003", "0.016", "0.54%", "0.9995", "Good"],
    ["MTR004", "0.017", "0.47%", "0.9995", "Good"],
    ["MTR005", "0.019", "0.45%", "0.9993", "Best MAPE (industrial)"]],
    LEFT_M, Inches(6.7), Inches(11.7), [Inches(0.7), Inches(0.6), Inches(0.6), Inches(0.6), Inches(2.0)], rh=Inches(0.17))

footer(s, 7)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — FORECAST ERROR (1 chart + 1 table)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header_bar(s, "Forecast Error — Night Low-Load Pattern")

sh(s, Inches(1.5), "FORECAST WINDOWS — Actual vs Predicted Around Top Errors")
chrt(s, "forecast_windows.png", LEFT_M, Inches(1.8), Inches(11.7), Inches(4.0))
obs(s, LEFT_M, Inches(5.9), Inches(11.7), Inches(0.4),
    "Pattern", "Highest % errors (68.9%, 39.9%) at night when load near zero. Absolute errors <0.2 kW — within meter accuracy tolerance.", ORANGE)

sh(s, Inches(6.4), "TOP 3 ERRORS — Root Cause")
tbl(s, [["Meter", "Actual", "Predicted", "Error", "Root Cause"],
    ["MTR005 12:30", "6.553 kW", "6.382 kW", "0.171 kW (2.6%)", "Industrial variability"],
    ["MTR002 23:30", "0.385 kW", "0.539 kW", "0.154 kW (39.9%)", "Night low-load"],
    ["MTR001 01:15", "0.184 kW", "0.311 kW", "0.127 kW (68.9%)", "Night low-load"]],
    LEFT_M, Inches(6.7), Inches(11.7), [Inches(1.0), Inches(0.9), Inches(0.9), Inches(1.1), Inches(1.3)], rh=Inches(0.17))

footer(s, 8)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — ANOMALY DETECTION (1 chart + 1 table)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header_bar(s, "Anomaly Detection — 6-Algorithm Ensemble")

sh(s, Inches(1.5), "ANOMALY RESULTS — 4,255 Anomalies (2.4%)")
chrt(s, "anomaly_detection_detail.png", LEFT_M, Inches(1.8), Inches(10.0), Inches(4.0))
obs(s, LEFT_M, Inches(5.9), Inches(10.0), Inches(0.4),
    "Key Finding", "MTR005 highest (1,284 = 3.66%). Night concentration 31.5%. Only 33 readings at severity 4+ (0.02%).", ORANGE)

sh(s, Inches(1.5), "ANOMALY SEVERITY — Action Required", Inches(10.8))
tbl(s, [["Score", "Count", "%", "Action"],
    ["0-1", "170,945", "97.6%", "Normal — no action"],
    ["2-3", "4,222", "2.4%", "Investigate"],
    ["4+", "33", "0.02%", "Immediate action"]],
    Inches(10.8), Inches(1.8), Inches(1.7), [Inches(0.4), Inches(0.5), Inches(0.4), Inches(0.9)], rh=Inches(0.22))

footer(s, 9)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — UC2→UC1 HANDOFF (1 flow + 1 table)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header_bar(s, "UC2 → UC1 Handoff: Timestamp Selection")

sh(s, Inches(1.5), "HANDOFF FLOW — UC2 Selects, UC1 Investigates")
steps = [("Load\nProfile", BLUE), ("Forecast", GREEN), ("Residual", ORANGE),
         ("Anomaly\nDetection", RED), ("Timestamp\nRanking", RED), ("UC1\nInvestigation", GREEN)]
for i, (step, c) in enumerate(steps):
    x = LEFT_M + Inches(i * 2.05)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(1.8), Inches(0.55), fill=c)
    txt(s, x, Inches(1.82), Inches(1.8), Inches(0.5), step, sz=10, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        txt(s, x + Inches(1.82), Inches(1.9), Inches(0.22), Inches(0.35), "->", sz=12, c=GRAY, a=PP_ALIGN.CENTER)

sh(s, Inches(2.6), "SELECTED TIMESTAMP — UC1 Investigation Target")
tbl(s, [["#", "Meter", "Timestamp", "Why Selected", "UC2 Finding"],
    ["1", "MTR005", "2025-11-06 13:30", "Highest peak (6.995 kW)", "1.64x 7-day avg"]],
    LEFT_M, Inches(3.3), Inches(11.7), [Inches(0.3), Inches(0.6), Inches(1.4), Inches(1.8), Inches(2.0)], rh=Inches(0.24))

obs(s, LEFT_M, Inches(4.5), Inches(11.7), Inches(0.4),
    "Rule", "UC1 ONLY investigates timestamps selected by UC2. No random timestamps. UC2 discovered → UC1 explains.", ORANGE)

footer(s, 10)

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — UC1 INVESTIGATION (1 chart + 1 table)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header_bar(s, "UC1 Investigation — MTR005 @ 2025-11-06 13:30")

sh(s, Inches(1.5), "ELECTRICAL TRENDS — Before → Event → After")
chrt(s, "event_electrical_trends.png", LEFT_M, Inches(1.8), Inches(9.0), Inches(4.0))
obs(s, LEFT_M, Inches(5.9), Inches(9.0), Inches(0.4),
    "Evidence", "Gradual ramp 4.95→6.99→5.56 kW. No sudden spike. V=224.1V stable. I=10.49A high. PF=0.975 excellent.", GREEN)

sh(s, Inches(1.5), "BEFORE → EVENT → AFTER", Inches(8.5))
tbl(s, [["Time", "Power", "Voltage", "Current", "PF"],
    ["15min before", "4.953 kW", "224.6 V", "7.92 A", "0.901"],
    ["EVENT", "6.995 kW", "224.1 V", "10.49 A", "0.975"],
    ["15min after", "6.244 kW", "233.5 V", "9.47 A", "0.945"]],
    Inches(8.5), Inches(1.8), Inches(4.0), [Inches(0.9), Inches(0.6), Inches(0.6), Inches(0.6), Inches(0.5)], rh=Inches(0.22))

obs(s, Inches(8.5), Inches(3.5), Inches(4.0), Inches(0.4),
    "Conclusion", "NORMAL — Industrial Activity. Gradual load increase, not sudden theft or fault.", GREEN)

footer(s, 11)

# ══════════════════════════════════════════════════════════════
# SLIDE 12 — HEALTH SCORE (1 formula + 1 table)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header_bar(s, "Meter Health Assessment — Component Breakdown")

sh(s, Inches(1.5), "HEALTH SCORE FORMULA")
add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, LEFT_M, Inches(1.8), Inches(11.7), Inches(0.5), fill=LIGHT_BLUE)
add_shape(s, MSO_SHAPE.RECTANGLE, LEFT_M, Inches(1.8), Pt(4), Inches(0.5), fill=BLUE)
txt(s, LEFT_M + Inches(0.15), Inches(1.83), Inches(11), Inches(0.2),
    "Overall = Comm(20%) + Voltage(20%) + PF(20%) + Anomaly(20%) + Tamper(10%) + Forecast(10%)", sz=12, c=WHITE, b=True)

sh(s, Inches(2.5), "CALCULATION TRAIL — Per Meter")
tbl(s, [["Component", "MTR001", "MTR002", "MTR003", "MTR004", "MTR005", "Formula"],
    ["Comm", "100", "100", "100", "100", "100", "100×valid/expected"],
    ["Voltage", "80.1", "80.0", "80.1", "79.9", "80.1", "100-std×5"],
    ["PF", "94.9", "94.9", "94.9", "94.9", "94.9", "avg×100"],
    ["Anomaly", "82.6", "90.7", "92.7", "91.6", "81.7", "100-rate×5"],
    ["Tamper*", "55", "55", "55", "55", "55", "9 DLMS events × 5 penalty"],
    ["Forecast", "85.8", "92.8", "94.6", "95.3", "95.5", "100-MAPE×5"],
    ["OVERALL", "85.6", "87.9", "88.5", "88.3", "86.4", "Weighted sum"],
    ["STATUS", "MONITOR", "MONITOR", "HEALTHY", "HEALTHY", "MONITOR", "≥88=HEALTHY"]],
    LEFT_M, Inches(2.8), Inches(11.7), [Inches(0.7), Inches(0.6), Inches(0.6), Inches(0.6), Inches(0.6), Inches(0.6), Inches(1.2)], rh=Inches(0.2))

obs(s, LEFT_M, Inches(5.2), Inches(11.7), Inches(0.4),
    "Key Finding", "MTR001=85.6 (MONITOR): Anomaly(82.6) + Tamper*(55, from 9 DLMS events/meter) drag score. MTR003=88.5 (HEALTHY): All components strong.", GREEN)

footer(s, 12)

# ══════════════════════════════════════════════════════════════
# SLIDE 13 — ROOT CAUSE (1 chart + 1 table)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header_bar(s, "Root Cause Analysis — Top 10 Peak Events")

sh(s, Inches(1.5), "ROOT CAUSE DISTRIBUTION — Top 10 Peak Events")
chrt(s, "root_cause_analysis.png", LEFT_M, Inches(1.8), Inches(10.0), Inches(4.0))
obs(s, LEFT_M, Inches(5.9), Inches(10.0), Inches(0.4),
    "Finding", "Industrial Activity (26.6%) is #1 cause across top 10 peaks. Commercial Behaviour (22.2%) is #2. Both are normal load patterns.", ORANGE)

sh(s, Inches(1.5), "TOP 3 CONFIRMED CAUSES", Inches(10.8))
tbl(s, [["Cause", "%", "Evidence"],
    ["Industrial Activity", "26.6%", "Weekday, afternoon, high current"],
    ["Commercial Behaviour", "22.2%", "Business hours, consistent load"],
    ["Current Increase", "19.3%", "I increases while V stable"],
    ["Other (voltage sag, comms)", "31.9%", "Remaining causes"]],
    Inches(10.8), Inches(1.8), Inches(1.7), [Inches(0.8), Inches(0.4), Inches(1.2)], rh=Inches(0.24))

sh(s, Inches(5.5), "HYPOTHESES EVALUATED AND RULED OUT")
obs(s, LEFT_M, Inches(5.8), Inches(11.7), Inches(0.5),
    "Ruled Out", "High Temperature (4.4%) and Residential Behaviour (4.4%) — insufficient evidence across top 10 peaks. DLMS events (270) are firmware artifacts, not theft.", GRAY)

footer(s, 13)

# ══════════════════════════════════════════════════════════════
# SLIDE 14 — CONCLUSIONS (KPIs + 1 table)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header_bar(s, "Conclusions & Recommendations")

kpis = [("FORECAST ACCURACY", "MAPE 0.72%", "R²=0.9997, MAE=0.017 kW", GREEN),
        ("ANOMALY DETECTION", "4,255 found", "2.4% of readings, 6-algorithm ensemble", ORANGE),
        ("HEALTH STATUS", "85.6 – 88.5", "No critical meters. MTR001 needs monitoring.", GREEN),
        ("THEFT DETECTION", "None Found", "0 bypass signatures. All DLMS artifacts.", GREEN)]
for i, (t, v, d, c) in enumerate(kpis):
    kpi(s, LEFT_M + Inches(i * 2.9), Inches(1.5), Inches(2.7), Inches(0.85), t, v, d, c)

sh(s, Inches(2.6), "ENGINEERING RECOMMENDATIONS")
tbl(s, [["Action", "Priority", "Evidence", "Timeline"],
    ["Monitor MTR001 voltage", "HIGH", "Health 85.6, voltage std 3.98V", "Immediate"],
    ["Deploy LightGBM pilot", "HIGH", "MAPE=0.72%, validated on 5 meters", "This quarter"],
    ["Monitor summer peaks", "MEDIUM", ">40C avg=4.10 kW, Jul highest", "Jun-Aug"],
    ["Scale to 5,000 meters", "MEDIUM", "Pilot validated, pipeline ready", "6 months"]],
    LEFT_M, Inches(2.9), Inches(11.7), [Inches(1.8), Inches(0.8), Inches(2.5), Inches(1.2)], rh=Inches(0.24))

footer(s, 14)

# ══════════════════════════════════════════════════════════════
# SLIDE 15 — CASE STUDY (1 chart + 1 table)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
header_bar(s, "Case Study: MTR005 — Complete Pipeline: Data to Decision")
txt(s, LEFT_M, Inches(1.1), Inches(10), Inches(0.3), "One meter, one complete engineering journey", sz=13, c=WHITE)

sh(s, Inches(1.5), "MTR005 — Year-Long Behaviour")
chrt(s, "mtr005_case_study.png", LEFT_M, Inches(1.8), Inches(11.7), Inches(4.0))
obs(s, LEFT_M, Inches(5.9), Inches(11.7), Inches(0.4),
    "Story", "Industrial consumer. Peak=6.995kW (Nov). Avg=4.26kW. Base=3.02kW. LF=0.609. Night/Day=0.64. 24h operation.", BLUE)

sh(s, Inches(5.9), "JOURNEY: Behaviour → Forecast → Anomaly → Investigation → Health")
tbl(s, [["Stage", "Finding", "Evidence"],
    ["1. Behaviour", "Industrial 24h", "Base=3.02kW, LF=0.609"],
    ["2. Monthly", "Nov peak=6.995kW", "Global maximum across all meters"],
    ["3. Forecast", "MAPE=0.45%", "Best among all meters"],
    ["4. Error", "Worst=0.171kW", "Industrial variability, not model failure"],
    ["5. Anomaly", "1,284 (3.66%)", "Highest count, expected for industrial"],
    ["6. Investigation", "NORMAL", "V=224.1V, I=10.49A, PF=0.975"],
    ["7. Health", "86.4 MONITOR", "Anomaly+Tamper components low"],
    ["8. Root Cause", "Industrial Activity", "Weekday, afternoon pattern"]],
    LEFT_M, Inches(6.2), Inches(11.7), [Inches(0.8), Inches(1.0), Inches(2.5)], rh=Inches(0.14))

footer(s, 15)

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
out = "/Users/yashitarora/Desktop/NES/Nes1/NES_UC1_UC2.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
